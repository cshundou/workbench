"""
PPT 文件生成服务（python-pptx）。
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.config import settings

logger = logging.getLogger(__name__)


def _safe_filename(name: str) -> str:
    """生成安全文件名。"""
    cleaned = re.sub(r'[/\\?%*:|"<>]', "_", name).strip("._ ")
    return cleaned[:80] or "presentation"


class PptGeneratorService:
    """演示文稿生成器。"""

    def __init__(self, base_dir: str | None = None) -> None:
        self.base_dir = Path(base_dir or settings.deliverables_dir)

    def resolve_session_dir(self, tenant_id: int, session_id: int) -> Path:
        """获取会话交付物目录并确保存在。"""
        path = self.base_dir / str(tenant_id) / str(session_id)
        path.mkdir(parents=True, exist_ok=True)
        return path

    def generate_pptx(
        self,
        tenant_id: int,
        session_id: int,
        outline: dict[str, Any],
        *,
        filename: str | None = None,
    ) -> dict[str, Any]:
        """
        根据大纲生成 PPTX 文件。

        Returns:
            {"filename", "file_path", "size", "slide_count"}
        """
        title = str(outline.get("title") or "演示文稿")
        slides_data = outline.get("slides") or []
        if not slides_data:
            slides_data = [{"title": title, "bullets": ["暂无内容"]}]

        session_dir = self.resolve_session_dir(tenant_id, session_id)
        safe_name = _safe_filename(filename or title)
        if not safe_name.lower().endswith(".pptx"):
            safe_name = f"{safe_name}.pptx"
        output_path = session_dir / safe_name

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 封面
        title_layout = prs.slide_layouts[0]
        cover = prs.slides.add_slide(title_layout)
        cover.shapes.title.text = title
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = "多 Agent 协同生成"

        content_layout = prs.slide_layouts[1]
        for slide_item in slides_data:
            if not isinstance(slide_item, dict):
                continue
            slide_title = str(slide_item.get("title") or "内容").strip()[:80]
            bullets_raw = slide_item.get("bullets") or []
            bullets = (
                [str(b).strip() for b in bullets_raw if str(b).strip()]
                if isinstance(bullets_raw, list)
                else [str(bullets_raw).strip()]
            )
            if not bullets:
                bullets = ["—"]

            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_title
            body = slide.placeholders[1].text_frame
            body.clear()
            for idx, bullet in enumerate(bullets[:8]):
                if idx == 0:
                    p = body.paragraphs[0]
                else:
                    p = body.add_paragraph()
                p.text = bullet[:300]
                p.level = 0
                p.font.size = Pt(18)

        prs.save(str(output_path))
        size = output_path.stat().st_size
        logger.info(
            "PPT 已生成 tenant=%s session=%s path=%s slides=%d",
            tenant_id,
            session_id,
            output_path,
            len(prs.slides),
        )
        return {
            "filename": safe_name,
            "file_path": str(output_path),
            "size": size,
            "slide_count": len(prs.slides),
        }

    def get_file_path(
        self,
        tenant_id: int,
        session_id: int,
        filename: str,
    ) -> Path | None:
        """校验并返回交付物文件路径（防目录穿越）。"""
        if ".." in filename or "/" in filename or "\\" in filename:
            return None
        path = self.resolve_session_dir(tenant_id, session_id) / filename
        if not path.is_file():
            return None
        try:
            path.resolve().relative_to(self.resolve_session_dir(tenant_id, session_id).resolve())
        except ValueError:
            return None
        return path


ppt_generator_service = PptGeneratorService()
