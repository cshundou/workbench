"""
PPT 生成引擎：模板、多页类型、图表、表格。
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings
from app.core.logging import get_logger
from app.services.ppt.schemas import PptOutline, SlideSpec
from app.services.ppt.templates import PptTemplate, get_template

logger = get_logger(__name__)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """十六进制颜色转 RGBColor。"""
    value = hex_color.lstrip("#")
    if len(value) != 6:
        value = "1F4E79"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _safe_filename(name: str) -> str:
    cleaned = re.sub(r'[/\\?%*:|"<>]', "_", name).strip("._ ")
    return cleaned[:80] or "presentation"


class PptGenerator:
    """演示文稿生成器（python-pptx）。"""

    def __init__(self, base_dir: str | None = None) -> None:
        self.base_dir = Path(base_dir or settings.deliverables_dir)

    def resolve_output_dir(self, tenant_id: int, session_id: int) -> Path:
        path = self.base_dir / str(tenant_id) / str(session_id)
        path.mkdir(parents=True, exist_ok=True)
        return path

    def generate(
        self,
        outline: dict[str, Any] | PptOutline,
        output_path: Path | str,
    ) -> dict[str, Any]:
        """
        根据大纲生成 PPTX。

        Args:
            outline: 结构化大纲 dict 或 PptOutline。
            output_path: 输出文件绝对路径。

        Returns:
            {"file_path", "filename", "size", "slide_count", "template_id"}
        """
        model = outline if isinstance(outline, PptOutline) else PptOutline.from_dict(outline)
        template = get_template(model.template_id)
        path = Path(output_path)
        if path.suffix.lower() != ".pptx":
            path = path.with_suffix(".pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        if not model.slides:
            self._add_cover(prs, model.title, model.subtitle or "多 Agent 协同生成", template)
            self._add_ending(prs, "谢谢", template)
        else:
            for spec in model.slides:
                slide_type = spec.slide_type or "content"
                if slide_type == "cover":
                    self._add_cover(prs, spec.title or model.title, spec.subtitle or model.subtitle, template)
                elif slide_type == "toc":
                    self._add_toc(prs, spec.title or "目录", spec.bullets, template)
                elif slide_type == "section":
                    self._add_section(prs, spec.title, spec.subtitle, template)
                elif slide_type == "ending":
                    self._add_ending(prs, spec.title or "谢谢聆听", template)
                else:
                    self._add_content_slide(prs, spec, template)

        prs.save(str(path))
        size = path.stat().st_size
        logger.info("PPT 生成完成 path=%s slides=%d template=%s", path, len(prs.slides), template.template_id)
        return {
            "file_path": str(path),
            "filename": path.name,
            "size": size,
            "slide_count": len(prs.slides),
            "template_id": template.template_id,
        }

    def generate_for_session(
        self,
        tenant_id: int,
        session_id: int,
        outline: dict[str, Any] | PptOutline,
        *,
        filename: str | None = None,
    ) -> dict[str, Any]:
        """写入租户/会话目录并生成 PPTX。"""
        model = outline if isinstance(outline, PptOutline) else PptOutline.from_dict(outline)
        out_dir = self.resolve_output_dir(tenant_id, session_id)
        safe = _safe_filename(filename or model.title)
        if not safe.lower().endswith(".pptx"):
            safe = f"{safe}.pptx"
        return self.generate(model, out_dir / safe)

    def get_file_path(self, tenant_id: int, session_id: int, filename: str) -> Path | None:
        """安全获取交付物路径。"""
        if ".." in filename or "/" in filename or "\\" in filename:
            return None
        base = self.resolve_output_dir(tenant_id, session_id)
        path = base / filename
        if not path.is_file():
            return None
        try:
            path.resolve().relative_to(base.resolve())
        except ValueError:
            return None
        return path

    def _apply_title_style(self, shape: Any, template: PptTemplate, *, size: int | None = None) -> None:
        if not shape or not shape.text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = template.title_font
            paragraph.font.size = Pt(size or template.title_size_pt)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _hex_to_rgb(template.primary_color)

    def _apply_body_style(self, text_frame: Any, template: PptTemplate) -> None:
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = template.body_font
            paragraph.font.size = Pt(template.body_size_pt)
            paragraph.font.color.rgb = _hex_to_rgb("333333")

    def _add_cover(
        self,
        prs: Presentation,
        title: str,
        subtitle: str,
        template: PptTemplate,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, template.background_color)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1.5))
        title_box.text_frame.text = title[:100]
        self._apply_title_style(title_box, template, size=40)
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
            sub_box.text_frame.text = subtitle[:200]
            for p in sub_box.text_frame.paragraphs:
                p.font.size = Pt(20)
                p.font.color.rgb = _hex_to_rgb(template.secondary_color)

    def _add_toc(
        self,
        prs: Presentation,
        title: str,
        items: list[str],
        template: PptTemplate,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, template.background_color)
        header = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
        header.text_frame.text = title
        self._apply_title_style(header, template, size=28)
        body = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10), Inches(5))
        tf = body.text_frame
        tf.clear()
        for idx, item in enumerate(items[:12]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"{idx + 1}. {item[:80]}"
            p.level = 0
        self._apply_body_style(tf, template)

    def _add_section(
        self,
        prs: Presentation,
        title: str,
        subtitle: str,
        template: PptTemplate,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, _strip_hash(template.secondary_color))
        box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.2))
        box.text_frame.text = title[:80]
        for p in box.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
            sub.text_frame.text = subtitle[:120]
            for p in sub.text_frame.paragraphs:
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(230, 230, 230)
                p.alignment = PP_ALIGN.CENTER

    def _add_ending(self, prs: Presentation, title: str, template: PptTemplate) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, template.primary_color)
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        box.text_frame.text = title[:50]
        for p in box.text_frame.paragraphs:
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        prs: Presentation,
        spec: SlideSpec,
        template: PptTemplate,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, template.background_color)
        header = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
        header.text_frame.text = (spec.title or "内容")[:80]
        self._apply_title_style(header, template, size=26)

        top = 1.4
        if spec.bullets or spec.paragraphs:
            body = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(5.5), Inches(5.5))
            tf = body.text_frame
            tf.clear()
            lines = spec.bullets or spec.paragraphs
            for idx, line in enumerate(lines[:10]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = str(line)[:300]
                p.level = 0
            self._apply_body_style(tf, template)

        if spec.chart and spec.chart.categories:
            self._add_chart(slide, spec.chart, template, left=Inches(6.5), top=Inches(1.6))

        if spec.table and spec.table.headers:
            self._add_table(slide, spec.table, top=Inches(4.5) if spec.chart else Inches(1.6))

    def _add_chart(
        self,
        slide: Any,
        chart_spec: Any,
        template: PptTemplate,
        *,
        left: Any,
        top: Any,
    ) -> None:
        chart_data = ChartData()
        chart_data.categories = chart_spec.categories[:12]
        for series in chart_spec.series[:5]:
            name = str(series.get("name", "系列"))
            values = series.get("values") or []
            chart_data.add_series(name, tuple(values[: len(chart_spec.categories)]))

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_type = chart_type_map.get(chart_spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        slide.shapes.add_chart(xl_type, left, top, Inches(6), Inches(4), chart_data)

    def _add_table(self, slide: Any, table_spec: Any, *, top: Any) -> None:
        rows = len(table_spec.rows) + 1
        cols = max(len(table_spec.headers), 1)
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), top, Inches(11.5), Inches(0.4 * rows))
        table = table_shape.table
        for col_idx, header in enumerate(table_spec.headers[:cols]):
            table.cell(0, col_idx).text = str(header)[:50]
        for row_idx, row in enumerate(table_spec.rows[:20]):
            for col_idx, cell in enumerate(row[:cols]):
                table.cell(row_idx + 1, col_idx).text = str(cell)[:80]

    @staticmethod
    def _set_slide_bg(slide: Any, color_hex: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(_strip_hash(color_hex))


def _strip_hash(color: str) -> str:
    return color.lstrip("#")


ppt_generator = PptGenerator()
